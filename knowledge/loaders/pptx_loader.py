from pptx import Presentation
from loaders.base import BaseLoader

class PPTXLoader(BaseLoader):
    def load(self, path: str) -> str:
        prs = Presentation(path)
        lines = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"## Slide {slide_idx}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        lines.append(text)

        return "\n\n".join(lines)
